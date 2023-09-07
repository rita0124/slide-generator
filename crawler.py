from pptx import Presentation

# 讀進樣板
f = open('./ppt/材料組進度報告 - 樣板.pptx', 'rb')
prs = Presentation(f)
f.close()

for template in prs.slide_layouts:
    # Insert this slide layout into the presentation
    slide = prs.slides.add_slide(template)
    # Find out shapes in this slide
    shape_count = 0
    for shape in slide.shapes:
        # Find out text in this shape
        if shape.has_text_frame:
            # Find out text in this text frame
            shape.text = str(shape_count)
        shape_count += 1

prs.save('./ppt/樣板標記.pptx')
