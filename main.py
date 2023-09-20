import re
import configparser
from pptx import Presentation
from datetime import date, datetime
from libs.kanban import Kanban
from libs.chart import Chart


class PPTBoss():

    def __init__(self) -> None:
        config = configparser.ConfigParser()
        config.read('config.ini')
        self.default = config['Default']
        self.kb = Kanban()
        self.prs = None
        self.chart = Chart()

    def read_from_template(self):
        with open('./ppt/材料組進度報告 - 樣板.pptx', 'rb') as f:
            self.prs = Presentation(f)
            f.close()

    def gen_prs_title_page(self):
        title_slide_layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(title_slide_layout)

        title = slide.shapes[1]
        title.text = "材料組進度報告"

        subtitle = slide.shapes[0]
        subtitle.text = f"大老闆 - {datetime.now().strftime('%Y-%m-%d')}"

    def gen_prs_content_pages(self):
        count = 1
        for team_id, team_name in self.kb.teams.items():
            # to_insert = 'N'
            # to_insert = input(f'加入組別: {team_name} ?(Y或y新增)')
            # if to_insert in 'Yy':
            if '材料' in team_name: # My team
                users = self.kb.get_users_in_team(team_id)
                projects = self.kb.get_projects_in_team(team_id)
                for project_id, project_name in projects.items():
                    ### 希望在每個專案的尾巴，都加上甘特圖，所以要在這專案迴圈開始的時候，記錄那些含有甘特圖要素的任務
                    # dict(Task = "工作1", Start = '2023-09-13', End = '2023-10-25', Assigned = "Rita" ）
                    gantt_tasks = []    ### 準備要被塞給 Chart().load_data(tasks) 的 tasks
                    to_insert = 'Y'
                    # to_insert = input(f'加入專案: {project_name} ?(Y?)')
                    if to_insert in 'YyYESyesYes':
                        # self.kb.get_tasks_in_project(project_id)
                        for user in users.keys():
                            tasks_details = self.kb.get_tasks_in_project_details(assignee_gid=user, project_gid=project_id)
                            # print(f'{user} ~~')
                            # self.kb.get_tasks_in_project_details(assignee_gid=user, project_gid=project_id)
                            # self.kb.clean_empty_values_in_my_tasks()
                            # for task_id in self.kb.my_tasks.keys():
                            if tasks_details is None:
                                continue
                            for task_id in tasks_details.keys():
                                # print(f"{self.kb.my_tasks[task_id]}")
                                # if self.kb.my_tasks[task_id]['name'] == '無名氏':
                                #     break
                                content_slide_layout = self.prs.slide_layouts[1]
                                slide = self.prs.slides.add_slide(content_slide_layout)
                                print(f'Add ... {count}')
                                count += 1

                                title = slide.shapes[0]
                                title.text = f"{tasks_details[task_id]['name']}"
                                # title.text = f"{self.kb.my_tasks[task_id]['name']} - {self.kb.my_tasks[task_id]['assignee']['name']}"

                                body_shape = slide.placeholders[1]
                                tf = body_shape.text_frame

                                # if 'start_on' not in tasks_details[task_id].keys():
                                if 'start_on' not in tasks_details[task_id].keys():
                                    start_on = '???'
                                    # start_on = '2000-01-01'
                                    #pass  #continue
                                # elif tasks_details[task_id]['start_on'] is None:
                                #    continue
                                else:
                                    # print(tasks_details[task_id]['start_on'])
                                    # start_on = tasks_details[task_id]['start_on'].strftime("%Y-%m-%d")
                                    start_on = tasks_details[task_id]['start_on']

                                # if 'due_on' not in tasks_details[task_id].keys():
                                if 'due_on' not in tasks_details[task_id].keys():
                                    due_on = '???'
                                    # due_on = '2050-01-01'
                                    pass  #continue
                                # elif tasks_details[task_id]['due_on'] is None:
                                    # continue
                                else:
                                    # due_on = tasks_details[task_id]['due_on'].strftime("%Y-%m-%d")
                                    due_on = tasks_details[task_id]['due_on']

                                assign_to = tasks_details[task_id]['assignee']['name']

                                is_completed = tasks_details[task_id]['completed']
                                tf.text = f"任務名稱：{tasks_details[task_id]['name']}\n"
                                tf.text += f'隸屬專案: {project_name}\n'
                                tf.text += f"指派給： {assign_to}\n"
                                tf.text += f"預定時間：{start_on} ~ {due_on}\n"
                                if is_completed:
                                    tf.text += f"狀態：{tasks_details[task_id]['completed_at'].strftime('%Y-%m-%d')} 完成"
                                elif start_on == '???' or due_on == '???':
                                    print('未知')
                                    print(tasks_details[task_id])
                                    tf.text += f"狀態： 未知"
                                elif tasks_details[task_id]['start_on'] > date.today():
                                    tf.text += f"狀態： 尚未開始"
                                else:
                                    tf.text += f"狀態： 進行中..."

                                if start_on != '???' and due_on != '???':
                                    gantt_tasks.append({
                                        'Task': tasks_details[task_id]['name'],
                                        'Start': start_on,
                                        'End': due_on,
                                        'Assigned': assign_to
                                    })
                    else:
                        print(f'略過專案: {project_name}')
                        continue
                    print(gantt_tasks)
                    self.chart.load_data(tasks=gantt_tasks)
                    self.chart.create_chart()
                    self.chart.save_chart(f'pics/gantt_{project_name}.png')
                    self.gen_gantt_pages(f'pics/gantt_{project_name}.png')
            else:
                print(f'略過組別: {team_name}')
                continue

    def gen_gantt_pages(self, chart_path=None):
        if chart_path is None:
            chart_path = 'pics/甘特圖.png'
        title_text = re.split('_|\.', chart_path)[1]
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        title = slide.shapes.title.text = f'{title_text}'
        # for shape in slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        placeholder = slide.placeholders[13]
        image = placeholder.insert_picture(chart_path)

    def save_file(self, filename=None):
        if filename is None:
            filename = f'./ppt/材料組進度報告_{datetime.today().strftime("%Y%m%d")}.pptx'
        self.prs.save(filename)


if __name__ == '__main__':
    pptboss = PPTBoss()
    pptboss.read_from_template()
    pptboss.gen_prs_title_page()
    pptboss.gen_prs_content_pages()
    pptboss.gen_gantt_pages()
    pptboss.save_file()
    print('done')